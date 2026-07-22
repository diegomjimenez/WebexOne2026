"""Enumerate slide layouts, their placeholders, and existing example slides in the
WebexOne 2026 template.

Run:  .venv\\Scripts\\python.exe scripts\\inspect_template.py

This is an analysis helper (OpenSpec task 2.1). It prints, for every slide master
layout, the layout name and each placeholder's idx / type / name, plus a summary of
the example slides shipped in the template (used to pick layouts and locate rich
motifs to duplicate).
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

TEMPLATE = Path(__file__).resolve().parent.parent / "templates" / "webexone-2026-light.pptx"


def _first_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    text = shape.text_frame.text.replace("\n", " / ").strip()
    return (text[:60] + "...") if len(text) > 60 else text


def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template not found: {TEMPLATE}", file=sys.stderr)
        return 1

    prs = Presentation(str(TEMPLATE))
    print(f"Opened template OK: {TEMPLATE.name}")
    print(f"Slide size: {prs.slide_width} x {prs.slide_height} EMU")
    print(f"Masters: {len(prs.slide_masters)}  Layouts (total): "
          f"{sum(len(m.slide_layouts) for m in prs.slide_masters)}  "
          f"Example slides: {len(prs.slides)}")

    print("\n==================== SLIDE LAYOUTS ====================")
    for mi, master in enumerate(prs.slide_masters):
        print(f"\n# Master {mi}: {master.name!r}")
        for li, layout in enumerate(master.slide_layouts):
            print(f"  [{mi}.{li}] layout name = {layout.name!r}")
            for ph in layout.placeholders:
                print(f"        ph idx={ph.placeholder_format.idx:<3} "
                      f"type={str(ph.placeholder_format.type):<22} "
                      f"name={ph.name!r} text={_first_text(ph)!r}")

    print("\n==================== EXAMPLE SLIDES ====================")
    for si, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name
        placeholders = [p.placeholder_format.idx for p in slide.placeholders]
        n_pics = sum(1 for s in slide.shapes if s.shape_type == 13)  # PICTURE
        print(f"  slide{si:<2} layout={layout_name!r:<40} "
              f"phs={placeholders} pics={n_pics}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
