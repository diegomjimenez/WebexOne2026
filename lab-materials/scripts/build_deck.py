"""Generate the WebexOne 2026 WxCC MCP lab deck from the official template.

Run (from lab-materials/):
    .venv\\Scripts\\python.exe scripts\\build_deck.py

Reads content/deck.yaml, opens templates/webexone-2026-light.pptx, and writes
build/WebexOne2026-WxCC-MCP-Lab.pptx. All styling (theme, fonts, colors) is
inherited from the template; this script only sets text and reuses layouts.

OpenSpec tasks 4.1 (generator), 4.2 (duplicate-and-retext helper), 4.3 (deterministic output).
"""

from __future__ import annotations

import copy
import re
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "templates" / "webexone-2026-light.pptx"
CONTENT = ROOT / "content" / "deck.yaml"
OUTPUT = ROOT / "build" / "WebexOne2026-WxCC-MCP-Lab.pptx"

# Logical role -> template layout name (see TEMPLATE-ANALYSIS.md, task 2.2).
LAYOUT_MAP = {
    "title": "WebexOne 2026 Title Slide 1",
    "contents": "Agenda 1",
    "section": "Section, Title Only 1",
    "concept": "Title, 1 Column with Bullets",
    "lab": "Title, 1 Column with Bullets",
    "statement": "Statement 1, Title, Subtitle",
    "resources": "1/2 Slide, Title, Body Copy, Graphic 1",
    "thankyou": "Thank you",
}

# 1-based index of the template example slide cloned for the metrics motif (task 2.3).
METRICS_SOURCE_SLIDE = 33
FOOTER_TYPE = 15  # PP_PLACEHOLDER.FOOTER

_ZW = re.compile(r"[\u200b\u200c\u200d\ufeff]")


def _norm(text: str) -> str:
    # Drop zero-width chars, then collapse all whitespace (incl. \x0b line breaks) to single spaces.
    return re.sub(r"\s+", " ", _ZW.sub("", text).replace("\xa0", " ")).strip()


def get_layout(prs, role: str):
    name = LAYOUT_MAP[role]
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f"Layout {name!r} (role {role!r}) not found in template")


def ph(slide, idx):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def set_lines(placeholder, lines, numbered=False):
    tf = placeholder.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = (f"{i + 1}. {line}") if numbered else line


def set_footer(slide, session_id):
    if not session_id:
        return
    for p in slide.placeholders:
        if p.placeholder_format.type == FOOTER_TYPE:
            p.text_frame.text = f"Session ID: {session_id}"


def clear_slides(prs):
    """Remove every example slide: drop the presentation->slide relationship (so the
    part becomes unreachable and is not re-serialized) and remove its sldId entry.
    Only dropping the sldId leaves orphaned parts that collide on save."""
    part = prs.part
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rid = sld_id.get(qn("r:id"))
        sld_id_lst.remove(sld_id)
        part.drop_rel(rid)


def duplicate_slide(prs, source_slide):
    """Deep-copy an example slide's shapes onto a new slide with the same layout.

    Valid for the metrics source slide, which has no pictures/relationships to remap.
    """
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in source_slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    return new_slide


def retext_by_map(slide, mapping):
    """Replace whole-paragraph text using a normalized-key mapping.

    Sets the paragraph's first run to the replacement and blanks the rest, so the
    template's run formatting (size/color) is preserved.
    """
    norm_map = {_norm(k): v for k, v in mapping.items()}
    hits = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            key = _norm(para.text)
            if key in norm_map and para.runs:
                para.runs[0].text = norm_map[key]
                for run in para.runs[1:]:
                    run.text = ""
                hits += 1
    return hits


def build_placeholder_slide(prs, role, spec, meta):
    slide = prs.slides.add_slide(get_layout(prs, role))

    if role == "title":
        slide.shapes.title.text = meta.get("title", "")
        if ph(slide, 13):
            ph(slide, 13).text_frame.text = meta.get("subtitle", "")
        if ph(slide, 11):
            ph(slide, 11).text_frame.text = meta.get("speaker1", "")
        if ph(slide, 14):
            ph(slide, 14).text_frame.text = meta.get("speaker2", "")
        if ph(slide, 12):
            ph(slide, 12).text_frame.text = meta.get("date", "")

    elif role == "contents":
        slide.shapes.title.text = spec.get("title", "Contents")
        if ph(slide, 20) and spec.get("numbers"):
            set_lines(ph(slide, 20), spec["numbers"])
        if ph(slide, 18) and spec.get("items"):
            set_lines(ph(slide, 18), spec["items"])

    elif role == "section":
        slide.shapes.title.text = spec.get("number", "")
        if ph(slide, 12):
            ph(slide, 12).text_frame.text = spec.get("title", "")

    elif role in ("concept", "lab"):
        title = spec.get("title", "")
        if role == "lab" and spec.get("mode"):
            title = f"{spec['mode']} — {title}"
        slide.shapes.title.text = title
        body = ph(slide, 12)
        if role == "lab":
            set_lines(body, spec.get("steps", []), numbered=True)
        else:
            lines = list(spec.get("bullets", []))
            if spec.get("reference"):
                lines.append(f"\u2192 {spec['reference']}")
            set_lines(body, lines)

    elif role == "statement":
        slide.shapes.title.text = spec.get("title", "")
        if ph(slide, 11):
            ph(slide, 11).text_frame.text = spec.get("subtitle", "")

    elif role == "resources":
        slide.shapes.title.text = spec.get("title", "Resources")
        if ph(slide, 11):
            set_lines(ph(slide, 11), spec.get("bullets", []))

    elif role == "thankyou":
        tf = slide.shapes.title.text_frame
        tf.text = spec.get("title", "Thank you")
        if spec.get("subtitle"):
            tf.add_paragraph().text = spec["subtitle"]

    else:
        raise ValueError(f"Unknown layout role: {role!r}")

    return slide


def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template not found: {TEMPLATE}", file=sys.stderr)
        return 1

    deck = yaml.safe_load(CONTENT.read_text(encoding="utf-8"))
    meta = deck.get("meta", {})
    session_id = meta.get("session_id")

    prs = Presentation(str(TEMPLATE))
    metrics_source = prs.slides[METRICS_SOURCE_SLIDE - 1]  # capture before clearing
    clear_slides(prs)

    built = 0
    for spec in deck["slides"]:
        role = spec["layout"]
        if role == "metrics_clone":
            slide = duplicate_slide(prs, metrics_source)
            mapping = dict(spec.get("replace", {}))
            if spec.get("title"):
                mapping["Growth metrics"] = spec["title"]
            hits = retext_by_map(slide, mapping)
            if hits < len(mapping):
                print(f"  ! metrics_clone: {hits}/{len(mapping)} replacements matched",
                      file=sys.stderr)
        else:
            slide = build_placeholder_slide(prs, role, spec, meta)
        set_footer(slide, session_id)
        built += 1

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT.relative_to(ROOT)} ({built} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
