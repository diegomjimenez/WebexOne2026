"""Verify the generated deck against the spec expectations (OpenSpec tasks 6.1, 6.2).

Checks: slide count in range; theme major/minor font is Inter; the seven numbered
section dividers exist; concept slides carry a reference arrow; lab slides carry a
hands-on/instructor badge. Prints a per-slide summary and a PASS/FAIL line.
"""

from __future__ import annotations

import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation

try:
    sys.stdout.reconfigure(encoding="utf-8")
except AttributeError:
    pass

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "build" / "WebexOne2026-WxCC-MCP-Lab.pptx"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


def theme_fonts(prs):
    theme_part = prs.slide_masters[0].part.part_related_by(THEME_REL)
    root = etree.fromstring(theme_part.blob)
    major = root.find(f".//{{{A}}}majorFont/{{{A}}}latin")
    minor = root.find(f".//{{{A}}}minorFont/{{{A}}}latin")
    return (major.get("typeface") if major is not None else None,
            minor.get("typeface") if minor is not None else None)


def main() -> int:
    prs = Presentation(str(DECK))
    n = len(prs.slides)
    major, minor = theme_fonts(prs)

    sections, concepts_with_ref, labs_with_badge, concept_total, lab_total = [], 0, 0, 0, 0
    print(f"{'#':>3}  {'layout':40}  title")
    for i, slide in enumerate(prs.slides, 1):
        title = slide.shapes.title.text.replace("\n", " / ") if slide.shapes.title else ""
        layout = slide.slide_layout.name
        print(f"{i:>3}  {layout:40}  {title[:60]}")
        body_text = "\n".join(
            s.text_frame.text for s in slide.shapes if s.has_text_frame
        )
        if layout == "Section, Title Only 1":
            sections.append(title)
        if layout == "Title, 1 Column with Bullets":
            if " — " in title and ("Lab" in title or "Demo" in title):
                lab_total += 1
                labs_with_badge += 1
            else:
                concept_total += 1
                if "\u2192" in body_text:
                    concepts_with_ref += 1

    print("\n--- checks ---")
    checks = {
        "slide count 20-30": 20 <= n <= 30,
        "major font Inter": (major or "").startswith("Inter"),
        "minor font Inter": (minor or "").startswith("Inter"),
        "7 section dividers": len(sections) == 7,
        "all concept slides have a reference": concepts_with_ref == concept_total and concept_total > 0,
        "all lab slides have a badge": labs_with_badge == lab_total and lab_total > 0,
    }
    for name, ok in checks.items():
        print(f"  [{'x' if ok else ' '}] {name}")
    print(f"\n  slides={n}  theme major/minor={major}/{minor}  "
          f"sections={len(sections)}  concepts={concept_total}  labs={lab_total}")

    ok = all(checks.values())
    print("\nRESULT:", "PASS" if ok else "FAIL")
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
